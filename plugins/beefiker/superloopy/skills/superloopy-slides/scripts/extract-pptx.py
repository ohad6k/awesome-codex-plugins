#!/usr/bin/env python3
"""
Extract all content from a PowerPoint file (.pptx).
Returns a JSON structure with slides, text, and images.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]

Requires: pip install python-pptx
"""

import json
import os
import sys
from pptx import Presentation

GROUP_SHAPE_TYPE = 6  # MSO_SHAPE_TYPE.GROUP
PICTURE_SHAPE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE


def walk_shape(shape, slide_num, slide_data, assets_dir, title_shape):
    """
    Extract text, tables, and images from one shape, recursing into groups.
    Table cells and text boxes nested in grouped shapes are real content —
    skipping them silently produces a converted deck missing text.
    """
    # Recurse into grouped shapes
    if shape.shape_type == GROUP_SHAPE_TYPE:
        for sub_shape in shape.shapes:
            walk_shape(sub_shape, slide_num, slide_data, assets_dir, title_shape)
        return

    # Extract table cell text (row by row, tab-separated)
    if getattr(shape, "has_table", False) and shape.has_table:
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if any(cell for row in rows for cell in row):
            slide_data["content"].append({"type": "table", "rows": rows})
        return

    # Extract text content
    if shape.has_text_frame and shape.text.strip():
        if title_shape is not None and shape == title_shape:
            slide_data["title"] = shape.text
        else:
            slide_data["content"].append({"type": "text", "content": shape.text})

    # Extract images (also found inside groups via the recursion above)
    if shape.shape_type == PICTURE_SHAPE_TYPE:
        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext
        image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
        image_path = os.path.join(assets_dir, image_name)

        with open(image_path, "wb") as f:
            f.write(image_bytes)

        slide_data["images"].append(
            {
                "path": f"assets/{image_name}",
                "width": shape.width,
                "height": shape.height,
            }
        )


def extract_pptx(file_path, output_dir="."):
    """
    Extract all content from a PowerPoint file.
    Returns a list of slide data dicts with text, tables, images, and notes.
    """
    prs = Presentation(file_path)
    slides_data = []

    # Create assets directory for extracted images
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        title_shape = slide.shapes.title  # None when the layout has no title
        for shape in slide.shapes:
            walk_shape(shape, slide_num, slide_data, assets_dir, title_shape)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data["notes"] = notes_frame.text

        slides_data.append(slide_data)

    return slides_data


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract-pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    slides = extract_pptx(input_file, output_dir)

    # Write extracted data as JSON
    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w") as f:
        json.dump(slides, f, indent=2)

    print(f"Extracted {len(slides)} slides to {output_path}")
    for s in slides:
        img_count = len(s["images"])
        print(f"  Slide {s['number']}: {s['title'] or '(no title)'} — {img_count} image(s)")
